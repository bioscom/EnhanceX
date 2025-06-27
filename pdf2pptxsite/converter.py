# converter.py
import fitz, pdfplumber, os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

def convert_pdf_to_pptx(pdf_path, output_pptx_path):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    temp_dir = "temp_img"
    os.makedirs(temp_dir, exist_ok=True)

    def fitz_to_inches(rect):
        return tuple(val * 0.01389 for val in rect)

    def create_textbox(slide, text, x, y, w, h, font_size=14, bold=False, color=(0, 0, 0)):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        p = box.text_frame.paragraphs[0]
        p.text = text.strip()
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        box.text_frame.word_wrap = True

    def insert_table(slide, table_data, x, y):
        rows, cols = len(table_data), len(table_data[0])
        table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(6), Inches(0.8 + 0.2 * rows)).table
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(table_data[r][c])
                cell.text_frame.fit_text(max_size=12)

    doc = fitz.open(pdf_path)
    plumber_doc = pdfplumber.open(pdf_path)

    for i in range(len(doc)):
        page = doc[i]
        ppage = plumber_doc.pages[i]
        slide = prs.slides.add_slide(blank)

        # Text blocks
        for block in page.get_text("dict")["blocks"]:
            if "lines" not in block: continue
            text, sizes = "", []
            for line in block["lines"]:
                for span in line["spans"]:
                    text += span["text"] + " "
                    sizes.append(span["size"])
            if text.strip():
                x, y, w, h = fitz_to_inches(block["bbox"])
                is_heading = min(sizes) >= 16 if sizes else False
                create_textbox(slide, text, x, y, w, h, min(sizes), is_heading)

        # Tables
        for table_data in ppage.extract_tables():
            if table_data and len(table_data) > 1:
                insert_table(slide, table_data, x=0.5, y=6.0)

        # Images
        for idx, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            ext = base_image["ext"]
            fname = os.path.join(temp_dir, f"{i}_{idx}.{ext}")
            with open(fname, "wb") as f:
                f.write(image_bytes)
            image = Image.open(fname)
            width, height = image.size
            slide.shapes.add_picture(fname, Inches(6), Inches(0.5),
                                     Inches(2.5), Inches(2.5 * height / width))

    prs.save(output_pptx_path)
    #print(f"✅ PowerPoint created: {output_pptx}")
    shutil.rmtree(temp_dir)